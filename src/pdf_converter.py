import os
from pdf2image import convert_from_path

# PDF → JPG
def pdf_para_jpg(caminho_pdf, pasta_saida=None, dpi=300):
    if pasta_saida is None:
        pasta_saida = os.path.dirname(caminho_pdf)
    os.makedirs(pasta_saida, exist_ok=True)

    paginas = convert_from_path(caminho_pdf, dpi=dpi)
    base = os.path.splitext(os.path.basename(caminho_pdf))[0]

    image_paths = []
    for i, pagina in enumerate(paginas, 1):
        nome_arquivo = os.path.join(pasta_saida, f'{base}_pagina_{i}.jpg')
        pagina.save(nome_arquivo, 'JPEG')
        image_paths.append(nome_arquivo)
    return image_paths

# PDF → PNG
def pdf_para_png(caminho_pdf, pasta_saida=None, dpi=300):
    if pasta_saida is None:
        pasta_saida = os.path.dirname(caminho_pdf)
    os.makedirs(pasta_saida, exist_ok=True)

    paginas = convert_from_path(caminho_pdf, dpi=dpi)
    base = os.path.splitext(os.path.basename(caminho_pdf))[0]

    image_paths = []
    for i, pagina in enumerate(paginas, 1):
        nome_arquivo = os.path.join(pasta_saida, f'{base}_pagina_{i}.png')
        pagina.save(nome_arquivo, 'PNG')
        image_paths.append(nome_arquivo)
    return image_paths

# PDF → TXT
def pdf_para_txt(caminho_pdf, pasta_saida=None):
    from pdfminer.high_level import extract_text
    if pasta_saida is None:
        pasta_saida = os.path.dirname(caminho_pdf)
    os.makedirs(pasta_saida, exist_ok=True)

    base = os.path.splitext(os.path.basename(caminho_pdf))[0]
    txt_path = os.path.join(pasta_saida, f'{base}.txt')
    texto = extract_text(caminho_pdf)
    with open(txt_path, 'w', encoding='utf-8') as f:
        f.write(texto)
    return txt_path

# PDF → DOCX
def pdf_para_docx(caminho_pdf, pasta_saida=None):
    from pdf2docx import Converter
    if pasta_saida is None:
        pasta_saida = os.path.dirname(caminho_pdf)
    os.makedirs(pasta_saida, exist_ok=True)

    base = os.path.splitext(os.path.basename(caminho_pdf))[0]
    docx_path = os.path.join(pasta_saida, f'{base}.docx')
    cv = Converter(caminho_pdf)
    cv.convert(docx_path, start=0)
    cv.close()
    return docx_path

# PDF → PowerPoint (PPTX)
def pdf_para_pptx(caminho_pdf, pasta_saida=None, dpi=200):
    from pptx import Presentation
    from pptx.util import Inches
    if pasta_saida is None:
        pasta_saida = os.path.dirname(caminho_pdf)
    os.makedirs(pasta_saida, exist_ok=True)

    base = os.path.splitext(os.path.basename(caminho_pdf))[0]
    pptx_path = os.path.join(pasta_saida, f'{base}.pptx')
    imagens = convert_from_path(caminho_pdf, dpi=dpi)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Slide em branco

    for idx, img in enumerate(imagens, 1):
        slide = prs.slides.add_slide(blank_slide_layout)
        img_path = os.path.join(pasta_saida, f'temp_{base}_{idx}.png')
        img.save(img_path, 'PNG')
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        os.remove(img_path)

    prs.save(pptx_path)
    return pptx_path